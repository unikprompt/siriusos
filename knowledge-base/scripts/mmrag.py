#!/usr/bin/env python3
"""
mmrag - Multimodal RAG Knowledge Base CLI

Ingest videos, images, audio, documents into a local ChromaDB vector database
using Gemini Embedding 2 for multimodal embeddings and Gemini Flash for
generating text descriptions of non-text media.

Usage:
    mmrag.py ingest <path> [<path>...] [--collection NAME]
    mmrag.py query <question> [--top-k N] [--threshold F] [--max-tokens N] [--collection NAME] [--json] [--full]
    mmrag.py status [--collection NAME]
    mmrag.py list [--collection NAME]
    mmrag.py collections
    mmrag.py delete <path> [--collection NAME]
    mmrag.py reset --confirm
"""

import argparse
import hashlib
import json
import mimetypes
import os
import subprocess
import sys
import time
from pathlib import Path

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
# cortextOS env-var overrides (set by kb-*.sh scripts)
MMRAG_DIR = Path(os.environ.get("MMRAG_DIR", str(Path.home() / ".mmrag")))
CONFIG_FILE = Path(os.environ.get("MMRAG_CONFIG", str(MMRAG_DIR / "config.json")))
CHROMADB_DIR = Path(os.environ.get("MMRAG_CHROMADB_DIR", str(MMRAG_DIR / "chromadb")))
MEDIA_DIR = MMRAG_DIR / "media"
LOG_DIR = MMRAG_DIR / "logs"

VIDEO_EXTS = {".mp4", ".mov", ".avi", ".mkv", ".webm"}
AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
DOC_EXTS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"}
TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".py", ".js", ".ts", ".go",
             ".rs", ".java", ".cpp", ".c", ".sh", ".yaml", ".yml", ".toml",
             ".html", ".css", ".sql", ".rb", ".swift", ".kt", ".r", ".lua"}

# Defaults
DEFAULT_TEXT_CHUNK_SIZE = 1500
DEFAULT_TEXT_CHUNK_OVERLAP = 200
DEFAULT_VIDEO_CHUNK_SECONDS = 60
DEFAULT_VIDEO_OVERLAP_SECONDS = 15
DEFAULT_AUDIO_CHUNK_SECONDS = 60
DEFAULT_AUDIO_OVERLAP_SECONDS = 10
DEFAULT_EMBEDDING_DIMENSIONS = 768
DEFAULT_SIMILARITY_THRESHOLD = 0.0  # return everything by default, let caller filter
DEFAULT_MAX_TOKENS = 0  # 0 = unlimited
DEFAULT_PREVIEW_CHARS = 300

# Pricing (per 1M tokens)
EMBEDDING_PRICE_PER_M = 0.20
FLASH_INPUT_PRICE_PER_M = 0.15
FLASH_OUTPUT_PRICE_PER_M = 0.60

USAGE_FILE = MMRAG_DIR / "usage.json"

# ---------------------------------------------------------------------------
# Usage Tracker
# ---------------------------------------------------------------------------
_tracker = None  # module-level, set by cmd_ingest/cmd_query


class UsageTracker:
    def __init__(self, operation="unknown"):
        self.session = {
            "embedding_tokens": 0,
            "embedding_calls": 0,
            "generation_input_tokens": 0,
            "generation_output_tokens": 0,
            "generation_calls": 0,
            "started_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
            "operation": operation,
        }

    def track_embedding(self, content):
        self.session["embedding_calls"] += 1
        if isinstance(content, str):
            self.session["embedding_tokens"] += int(len(content.split()) * 1.3)
        elif isinstance(content, list):
            for part in content:
                if isinstance(part, str):
                    self.session["embedding_tokens"] += int(len(part.split()) * 1.3)
                else:
                    try:
                        self.session["embedding_tokens"] += max(256, len(part.data) // 4)
                    except Exception:
                        self.session["embedding_tokens"] += 256

    def track_generation(self, response):
        self.session["generation_calls"] += 1
        um = getattr(response, "usage_metadata", None)
        if um:
            self.session["generation_input_tokens"] += getattr(um, "prompt_token_count", 0) or 0
            self.session["generation_output_tokens"] += getattr(um, "candidates_token_count", 0) or 0

    def cost(self):
        emb = (self.session["embedding_tokens"] / 1_000_000) * EMBEDDING_PRICE_PER_M
        gen_in = (self.session["generation_input_tokens"] / 1_000_000) * FLASH_INPUT_PRICE_PER_M
        gen_out = (self.session["generation_output_tokens"] / 1_000_000) * FLASH_OUTPUT_PRICE_PER_M
        return {
            "embedding": round(emb, 6),
            "generation_input": round(gen_in, 6),
            "generation_output": round(gen_out, 6),
            "total": round(emb + gen_in + gen_out, 6),
        }

    def persist(self):
        MMRAG_DIR.mkdir(parents=True, exist_ok=True)
        self.session["finished_at"] = time.strftime("%Y-%m-%dT%H:%M:%S")
        self.session["cost"] = self.cost()

        data = {"cumulative": {}, "sessions": []}
        if USAGE_FILE.exists():
            try:
                with open(USAGE_FILE) as f:
                    data = json.load(f)
            except (json.JSONDecodeError, KeyError):
                data = {"cumulative": {}, "sessions": []}

        data.setdefault("sessions", []).append(self.session)

        c = data.get("cumulative", {})
        for key in ["embedding_tokens", "embedding_calls",
                     "generation_input_tokens", "generation_output_tokens",
                     "generation_calls"]:
            c[key] = c.get(key, 0) + self.session[key]

        c["total_cost"] = round(sum(
            s.get("cost", {}).get("total", 0) for s in data["sessions"]
        ), 6)
        data["cumulative"] = c

        with open(USAGE_FILE, "w") as f:
            json.dump(data, f, indent=2)

    def summary_line(self):
        c = self.cost()
        return (f"  Tokens: {self.session['embedding_tokens']:,} embedding, "
                f"{self.session['generation_input_tokens']:,} gen-input, "
                f"{self.session['generation_output_tokens']:,} gen-output | "
                f"Cost: ${c['total']:.4f}")


# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------
def load_config():
    if not CONFIG_FILE.exists():
        print("ERROR: Config not found. Run setup first:")
        print(f"  bash {Path(__file__).parent / 'setup.sh'}")
        sys.exit(1)
    with open(CONFIG_FILE) as f:
        return json.load(f)


def get_api_key(config):
    key = os.environ.get("GEMINI_API_KEY") or config.get("gemini_api_key")
    if not key:
        print("ERROR: No Gemini API key. Set GEMINI_API_KEY or run setup.")
        sys.exit(1)
    return key

# ---------------------------------------------------------------------------
# Gemini clients
# ---------------------------------------------------------------------------
def get_genai_client(api_key):
    from google import genai
    return genai.Client(api_key=api_key)


def embed_content(client, config, content, task_type="RETRIEVAL_DOCUMENT"):
    """Embed content using Gemini Embedding 2. Content can be text string or list of Parts."""
    from google.genai import types
    result = client.models.embed_content(
        model=config.get("embedding_model", "gemini-embedding-2-preview"),
        contents=content,
        config=types.EmbedContentConfig(
            output_dimensionality=config.get("embedding_dimensions", DEFAULT_EMBEDDING_DIMENSIONS),
            task_type=task_type,
        ),
    )
    if _tracker:
        _tracker.track_embedding(content)
    return result.embeddings[0].values


def embed_multimodal(client, config, description_text, media_bytes, mime_type):
    """
    Option B embedding: combine text description + raw media into one embedding.
    This captures both semantic text meaning AND visual/audio content.
    """
    from google.genai import types
    contents = [
        description_text,
        types.Part.from_bytes(data=media_bytes, mime_type=mime_type),
    ]
    return embed_content(client, config, contents)


def embed_query(client, config, query_text):
    """Embed a query string for retrieval."""
    return embed_content(client, config, query_text, task_type="RETRIEVAL_QUERY")


def describe_media(client, config, file_path, media_type="video"):
    """Use Gemini Flash to generate a text description of media."""
    from google.genai import types

    mime = mimetypes.guess_type(str(file_path))[0] or "application/octet-stream"
    with open(file_path, "rb") as f:
        data = f.read()

    prompts = {
        "video": (
            "Provide a detailed description of this video. Include:\n"
            "1. What is being shown/demonstrated\n"
            "2. Any text visible on screen\n"
            "3. Key concepts or topics discussed\n"
            "4. A transcript of any spoken words\n"
            "5. Step-by-step actions if it's a tutorial\n"
            "Be thorough - this description will be used for search and retrieval."
        ),
        "image": (
            "Describe this image in detail. Include:\n"
            "1. What is shown in the image\n"
            "2. Any text visible in the image\n"
            "3. Key concepts or topics depicted\n"
            "4. Colors, layout, and composition\n"
            "Be thorough - this description will be used for search and retrieval."
        ),
        "audio": (
            "Transcribe and describe this audio. Include:\n"
            "1. A full transcript of spoken words\n"
            "2. Description of any sounds or music\n"
            "3. Key topics discussed\n"
            "4. Speaker identification if possible\n"
            "Be thorough - this description will be used for search and retrieval."
        ),
    }

    response = client.models.generate_content(
        model=config.get("gemini_model", "gemini-2.5-flash"),
        contents=[
            types.Part.from_bytes(data=data, mime_type=mime),
            prompts.get(media_type, prompts["video"]),
        ],
    )
    if _tracker:
        _tracker.track_generation(response)
    return response.text, data, mime

# ---------------------------------------------------------------------------
# ChromaDB
# ---------------------------------------------------------------------------
def get_chroma_collection(collection_name="default"):
    import chromadb
    client = chromadb.PersistentClient(path=str(CHROMADB_DIR))
    return client.get_or_create_collection(
        name=collection_name,
        metadata={"hnsw:space": "cosine"},
    )


def get_chroma_client():
    import chromadb
    return chromadb.PersistentClient(path=str(CHROMADB_DIR))

# ---------------------------------------------------------------------------
# Text chunking
# ---------------------------------------------------------------------------
def chunk_text(text, chunk_size=DEFAULT_TEXT_CHUNK_SIZE, overlap=DEFAULT_TEXT_CHUNK_OVERLAP):
    """Split text into overlapping chunks, preferring paragraph/section boundaries."""
    if len(text) <= chunk_size:
        return [text] if text.strip() else []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size

        # Try to break at a paragraph boundary (double newline)
        if end < len(text):
            # Look backwards from end for a good break point
            search_zone = text[max(start + chunk_size // 2, start):end]
            # Prefer double newline (paragraph break)
            para_break = search_zone.rfind("\n\n")
            if para_break != -1:
                end = max(start + chunk_size // 2, start) + para_break + 2
            else:
                # Fall back to single newline
                line_break = search_zone.rfind("\n")
                if line_break != -1:
                    end = max(start + chunk_size // 2, start) + line_break + 1

        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
        if start >= len(text):
            break

    return chunks

# ---------------------------------------------------------------------------
# Video chunking with FFmpeg
# ---------------------------------------------------------------------------
def get_media_duration(file_path):
    """Get duration of a media file in seconds. Returns 0 if unreadable."""
    try:
        result = subprocess.run(
            ["ffprobe", "-v", "quiet", "-show_entries", "format=duration",
             "-of", "csv=p=0", str(file_path)],
            capture_output=True, text=True,
        )
        val = result.stdout.strip()
        return float(val) if val else 0.0
    except (ValueError, subprocess.SubprocessError):
        return 0.0


def chunk_video(video_path, chunk_seconds=DEFAULT_VIDEO_CHUNK_SECONDS,
                overlap_seconds=DEFAULT_VIDEO_OVERLAP_SECONDS):
    """Split video into overlapping chunks using FFmpeg."""
    video_path = Path(video_path)
    output_dir = MEDIA_DIR / video_path.stem
    output_dir.mkdir(parents=True, exist_ok=True)

    duration = get_media_duration(video_path)

    chunks = []
    start = 0
    idx = 0
    step = chunk_seconds - overlap_seconds

    while start < duration:
        end = min(start + chunk_seconds, duration)
        # Skip tiny trailing chunks (< 5 seconds)
        if end - start < 5 and idx > 0:
            break

        output_file = output_dir / f"chunk_{idx:04d}.mp4"

        if not output_file.exists():
            subprocess.run(
                ["ffmpeg", "-y", "-i", str(video_path),
                 "-ss", str(start), "-t", str(end - start),
                 "-c", "copy", "-avoid_negative_ts", "1",
                 str(output_file)],
                capture_output=True,
            )

        chunks.append({
            "path": str(output_file),
            "start": start,
            "end": end,
            "index": idx,
        })

        start += step
        idx += 1

    return chunks


def chunk_audio(audio_path, chunk_seconds=DEFAULT_AUDIO_CHUNK_SECONDS,
                overlap_seconds=DEFAULT_AUDIO_OVERLAP_SECONDS):
    """Split audio into overlapping chunks using FFmpeg."""
    audio_path = Path(audio_path)
    output_dir = MEDIA_DIR / audio_path.stem
    output_dir.mkdir(parents=True, exist_ok=True)

    duration = get_media_duration(audio_path)

    ext = audio_path.suffix
    chunks = []
    start = 0
    idx = 0
    step = chunk_seconds - overlap_seconds

    while start < duration:
        end = min(start + chunk_seconds, duration)
        if end - start < 3 and idx > 0:
            break

        output_file = output_dir / f"chunk_{idx:04d}{ext}"

        if not output_file.exists():
            subprocess.run(
                ["ffmpeg", "-y", "-i", str(audio_path),
                 "-ss", str(start), "-t", str(end - start),
                 "-c", "copy", str(output_file)],
                capture_output=True,
            )

        chunks.append({
            "path": str(output_file),
            "start": start,
            "end": end,
            "index": idx,
        })

        start += step
        idx += 1

    return chunks

# ---------------------------------------------------------------------------
# File ID helper
# ---------------------------------------------------------------------------
def file_id(path, chunk_idx=None):
    """Generate a stable ID for a file or chunk."""
    h = hashlib.md5(str(path).encode()).hexdigest()[:12]
    if chunk_idx is not None:
        return f"{h}_chunk{chunk_idx}"
    return h

# ---------------------------------------------------------------------------
# Ingest logic
# ---------------------------------------------------------------------------
def already_exists(collection, doc_id):
    """Check if a document ID already exists in the collection. Respects --force flag."""
    if args_force:
        return False
    existing = collection.get(ids=[doc_id])
    return bool(existing and existing["ids"])


def ingest_text_file(client, config, collection, file_path):
    """Ingest a text-based file."""
    file_path = Path(file_path)
    text = file_path.read_text(errors="replace")
    if not text.strip():
        print(f"  SKIP (empty): {file_path}")
        return 0

    chunks = chunk_text(
        text,
        chunk_size=config.get("text_chunk_size", DEFAULT_TEXT_CHUNK_SIZE),
        overlap=config.get("text_chunk_overlap", DEFAULT_TEXT_CHUNK_OVERLAP),
    )

    count = 0
    for i, chunk in enumerate(chunks):
        doc_id = file_id(file_path, i)
        if already_exists(collection, doc_id):
            continue

        embedding = embed_content(client, config, chunk)
        collection.upsert(
            ids=[doc_id],
            embeddings=[embedding],
            documents=[chunk],
            metadatas=[{
                "source": str(file_path.resolve()),
                "type": "text",
                "chunk_index": i,
                "total_chunks": len(chunks),
                "filename": file_path.name,
                "file_ext": file_path.suffix.lower(),
                "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
            }],
        )
        count += 1

    return count


def ingest_image(client, config, collection, file_path):
    """Ingest an image: Gemini Flash describes it, then embed description + raw image together."""
    file_path = Path(file_path)
    doc_id = file_id(file_path)

    if already_exists(collection, doc_id):
        print(f"  SKIP (exists): {file_path}")
        return 0

    print(f"  Generating description for {file_path.name}...")
    description, media_bytes, mime = describe_media(client, config, file_path, "image")

    # Option B: embed text description + raw image together
    try:
        embedding = embed_multimodal(client, config, description, media_bytes, mime)
    except Exception:
        # Fallback to text-only embedding if multimodal fails (e.g., file too large)
        embedding = embed_content(client, config, description)

    collection.upsert(
        ids=[doc_id],
        embeddings=[embedding],
        documents=[description],
        metadatas=[{
            "source": str(file_path.resolve()),
            "type": "image",
            "filename": file_path.name,
            "file_ext": file_path.suffix.lower(),
            "mime_type": mime,
            "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
        }],
    )
    return 1


def extract_audio_from_video(video_path, output_path):
    """Extract audio track from a video file as mp3."""
    subprocess.run(
        ["ffmpeg", "-y", "-i", str(video_path),
         "-vn", "-acodec", "libmp3lame", "-q:a", "4",
         str(output_path)],
        capture_output=True,
    )
    return Path(output_path).exists()


def ingest_video(client, config, collection, file_path):
    """Ingest a video: chunk it, describe each chunk, embed.

    For large videos (chunks > 20MB), falls back to audio-only extraction
    since the video bytes would be too large for the embedding API.
    For small chunks, uses full multimodal embedding (description + video).
    """
    file_path = Path(file_path)
    size_mb = file_path.stat().st_size / (1024 * 1024)
    duration = get_media_duration(file_path)

    if duration <= 0:
        print(f"  SKIP (unreadable/zero duration): {file_path}")
        return 0

    print(f"  Video: {file_path.name} ({size_mb:.0f}MB, {duration:.0f}s)")

    # Chunk the video
    chunk_secs = config.get("video_chunk_seconds", DEFAULT_VIDEO_CHUNK_SECONDS)
    overlap_secs = config.get("video_overlap_seconds", DEFAULT_VIDEO_OVERLAP_SECONDS)

    print(f"  Chunking into {chunk_secs}s segments with {overlap_secs}s overlap...")
    chunks = chunk_video(file_path, chunk_seconds=chunk_secs, overlap_seconds=overlap_secs)
    total_chunks = len(chunks)
    print(f"  Created {total_chunks} chunks")

    count = 0
    for chunk in chunks:
        doc_id = file_id(file_path, chunk["index"])
        if already_exists(collection, doc_id):
            continue

        chunk_path = Path(chunk["path"])
        chunk_size_mb = chunk_path.stat().st_size / (1024 * 1024) if chunk_path.exists() else 0

        print(f"  Chunk {chunk['index'] + 1}/{total_chunks} "
              f"({chunk['start']:.0f}s-{chunk['end']:.0f}s, {chunk_size_mb:.1f}MB)")

        description = None
        media_bytes = None
        mime = None

        # Strategy: try video description first, fall back to audio-only for large chunks
        if chunk_size_mb <= 20:
            # Small enough for full video analysis
            try:
                description, media_bytes, mime = describe_media(client, config, chunk_path, "video")
                print(f"    Described via video")
            except Exception as e:
                print(f"    Video description failed ({e}), trying audio...")

        if description is None:
            # Large chunk or video failed: extract audio and describe that
            audio_path = chunk_path.with_suffix(".mp3")
            if not audio_path.exists():
                print(f"    Extracting audio track...")
                extract_audio_from_video(chunk_path, audio_path)

            if audio_path.exists() and audio_path.stat().st_size > 0:
                try:
                    description, media_bytes, mime = describe_media(client, config, audio_path, "audio")
                    # Prefix so the agent knows this came from a video's audio
                    description = (
                        f"[Audio extracted from video: {file_path.name}, "
                        f"{chunk['start']:.0f}s-{chunk['end']:.0f}s]\n\n{description}"
                    )
                    print(f"    Described via audio extraction")
                except Exception as e:
                    print(f"    Audio description also failed: {e}")

        if description is None:
            description = (
                f"Video chunk from {file_path.name}, "
                f"{chunk['start']:.0f}s to {chunk['end']:.0f}s. "
                f"(Description unavailable - file may be too large or corrupted)"
            )

        # Embed: try multimodal if we have small media bytes, else text-only
        if media_bytes and mime and len(media_bytes) < 20 * 1024 * 1024:
            try:
                embedding = embed_multimodal(client, config, description, media_bytes, mime)
            except Exception:
                embedding = embed_content(client, config, description)
        else:
            embedding = embed_content(client, config, description)

        collection.upsert(
            ids=[doc_id],
            embeddings=[embedding],
            documents=[description],
            metadatas=[{
                "source": str(file_path.resolve()),
                "type": "video_chunk",
                "chunk_index": chunk["index"],
                "total_chunks": total_chunks,
                "chunk_start_seconds": chunk["start"],
                "chunk_end_seconds": chunk["end"],
                "chunk_path": str(chunk_path),
                "filename": file_path.name,
                "file_ext": file_path.suffix.lower(),
                "duration_seconds": duration,
                "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
            }],
        )
        count += 1

    return count


def ingest_audio(client, config, collection, file_path):
    """Ingest audio: chunk if needed, describe, embed description + audio together."""
    file_path = Path(file_path)
    duration = get_media_duration(file_path)
    if duration <= 0:
        print(f"  SKIP (unreadable/zero duration): {file_path}")
        return 0
    max_chunk = config.get("audio_chunk_seconds", DEFAULT_AUDIO_CHUNK_SECONDS)

    if duration <= max_chunk:
        # Short enough to process as one piece
        doc_id = file_id(file_path)
        if already_exists(collection, doc_id):
            print(f"  SKIP (exists): {file_path}")
            return 0

        print(f"  Transcribing {file_path.name}...")
        description, media_bytes, mime = describe_media(client, config, file_path, "audio")

        try:
            embedding = embed_multimodal(client, config, description, media_bytes, mime)
        except Exception:
            embedding = embed_content(client, config, description)

        collection.upsert(
            ids=[doc_id],
            embeddings=[embedding],
            documents=[description],
            metadatas=[{
                "source": str(file_path.resolve()),
                "type": "audio",
                "filename": file_path.name,
                "file_ext": file_path.suffix.lower(),
                "duration_seconds": duration,
                "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
            }],
        )
        return 1
    else:
        # Chunk the audio
        print(f"  Chunking audio: {file_path.name} ({duration:.0f}s)...")
        overlap = config.get("audio_overlap_seconds", DEFAULT_AUDIO_OVERLAP_SECONDS)
        chunks = chunk_audio(file_path, chunk_seconds=max_chunk, overlap_seconds=overlap)
        total_chunks = len(chunks)
        count = 0

        for chunk in chunks:
            doc_id = file_id(file_path, chunk["index"])
            if already_exists(collection, doc_id):
                continue

            print(f"  Transcribing chunk {chunk['index'] + 1}/{total_chunks}...")
            try:
                description, media_bytes, mime = describe_media(client, config, chunk["path"], "audio")
            except Exception as e:
                print(f"  WARNING: Failed to transcribe chunk {chunk['index']}: {e}")
                description = f"Audio chunk from {file_path.name}, {chunk['start']:.0f}s to {chunk['end']:.0f}s"
                media_bytes = None
                mime = None

            if media_bytes and mime:
                try:
                    embedding = embed_multimodal(client, config, description, media_bytes, mime)
                except Exception:
                    embedding = embed_content(client, config, description)
            else:
                embedding = embed_content(client, config, description)

            collection.upsert(
                ids=[doc_id],
                embeddings=[embedding],
                documents=[description],
                metadatas=[{
                    "source": str(file_path.resolve()),
                    "type": "audio_chunk",
                    "chunk_index": chunk["index"],
                    "total_chunks": total_chunks,
                    "chunk_start_seconds": chunk["start"],
                    "chunk_end_seconds": chunk["end"],
                    "chunk_path": chunk["path"],
                    "filename": file_path.name,
                    "file_ext": file_path.suffix.lower(),
                    "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
                }],
            )
            count += 1
        return count


def ingest_pdf(client, config, collection, file_path):
    """Ingest a PDF page-by-page using Gemini to extract content including visual elements."""
    file_path = Path(file_path)
    from google.genai import types

    with open(file_path, "rb") as f:
        data = f.read()

    # Estimate page count (rough: ~3KB per page for typical PDFs, but varies wildly)
    # We'll ask Gemini to process the whole thing and get structured output
    # For PDFs > 6 pages, we chunk by asking for specific page ranges

    print(f"  Analyzing PDF: {file_path.name}...")

    # Gemini Flash returns 503 UNAVAILABLE during high-demand windows.
    # Without retries, a single 503 kills the ingest. Retry up to 3 times
    # with exponential backoff (5s, 15s, 45s). Re-raise on the last failure.
    # Predicate uses google.genai.errors.APIError's structured .code (HTTP int)
    # and .status (gRPC-style text). Substring matching on str(e) was rejected
    # because it false-positived non-transient errors whose body text
    # incidentally contained "500"/"503" (e.g. a 403 mentioning a resource id
    # with "503" in it would have wrongly triggered the retry loop).
    from google.genai import errors as _genai_errors
    TRANSIENT_HTTP_CODES = {429, 500, 503}
    TRANSIENT_STATUS_NAMES = {"UNAVAILABLE", "RESOURCE_EXHAUSTED"}
    extraction_prompt = (
        "Extract ALL content from this PDF. For each page, include:\n"
        "1. Page number\n"
        "2. All text content (headings, body, lists, footnotes)\n"
        "3. Description of any images, charts, diagrams, or tables\n"
        "4. Key concepts and topics on that page\n"
        "Separate each page's content with '=== PAGE N ===' markers.\n"
        "Be thorough - this will be used for search and retrieval."
    )
    response = None
    last_err = None
    for attempt, backoff in enumerate([5, 15, 45], start=1):
        try:
            response = client.models.generate_content(
                model=config.get("gemini_model", "gemini-2.5-flash"),
                contents=[
                    types.Part.from_bytes(data=data, mime_type="application/pdf"),
                    extraction_prompt,
                ],
            )
            break
        except _genai_errors.APIError as e:
            last_err = e
            # Structured retry predicate: only retry on real transient
            # SDK-level conditions. Auth / 4xx config errors fail fast even
            # when their response body text incidentally contains digits like
            # "503" — this was the false-positive class flagged in PR review.
            is_transient = (e.code in TRANSIENT_HTTP_CODES) or (e.status in TRANSIENT_STATUS_NAMES)
            if not is_transient:
                raise
            if attempt < 3:
                print(f"    Transient error (HTTP {e.code} {e.status or ''}); retrying in {backoff}s (attempt {attempt}/3)")
                time.sleep(backoff)
            else:
                print(f"    Exhausted retries on transient error: HTTP {e.code} {e.status or ''}")
    if response is None:
        raise last_err if last_err else RuntimeError("PDF ingest failed without exception captured")
    if _tracker:
        _tracker.track_generation(response)
    text = response.text

    # Split by page markers if present, otherwise chunk normally
    pages = []
    if "=== PAGE" in text:
        import re
        page_splits = re.split(r'===\s*PAGE\s*\d+\s*===', text)
        pages = [p.strip() for p in page_splits if p.strip()]
    else:
        # No page markers - chunk as text
        pages = chunk_text(
            text,
            chunk_size=config.get("text_chunk_size", DEFAULT_TEXT_CHUNK_SIZE),
            overlap=config.get("text_chunk_overlap", DEFAULT_TEXT_CHUNK_OVERLAP),
        )

    count = 0
    for i, page_content in enumerate(pages):
        if not page_content.strip():
            continue
        doc_id = file_id(file_path, i)
        if already_exists(collection, doc_id):
            continue

        embedding = embed_content(client, config, page_content)
        collection.upsert(
            ids=[doc_id],
            embeddings=[embedding],
            documents=[page_content],
            metadatas=[{
                "source": str(file_path.resolve()),
                "type": "pdf_page",
                "chunk_index": i,
                "total_chunks": len(pages),
                "page_number": i + 1,
                "filename": file_path.name,
                "file_ext": ".pdf",
                "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
            }],
        )
        count += 1
    return count


def extract_docx_text(file_path):
    """Extract text from .docx using python-docx."""
    from docx import Document
    doc = Document(str(file_path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            style_name = para.style.name if para.style else ""
            if style_name.startswith("Heading"):
                level = style_name.replace("Heading ", "").strip()
                prefix = "#" * (int(level) if level.isdigit() else 1)
                parts.append(f"{prefix} {para.text}")
            else:
                parts.append(para.text)
    # Also extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n".join(rows))
    return "\n\n".join(parts)


def extract_pptx_text(file_path):
    """Extract text from .pptx using python-pptx."""
    from pptx import Presentation
    prs = Presentation(str(file_path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = [f"=== SLIDE {i+1} ==="]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        # Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"Speaker Notes: {notes}")
        slides.append("\n".join(texts))
    return "\n\n".join(slides)


def extract_xlsx_text(file_path):
    """Extract text from .xlsx using openpyxl."""
    from openpyxl import load_workbook
    wb = load_workbook(str(file_path), data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"=== SHEET: {sheet_name} ===\n" + "\n".join(rows[:200]))  # cap at 200 rows
    return "\n\n".join(parts)


def ingest_office_doc(client, config, collection, file_path):
    """Ingest Office documents (.docx, .pptx, .xlsx) by extracting text locally."""
    file_path = Path(file_path)
    ext = file_path.suffix.lower()

    print(f"  Extracting content from {file_path.name}...")

    try:
        if ext in (".docx", ".doc"):
            text = extract_docx_text(file_path)
            type_name = "docx"
        elif ext in (".pptx", ".ppt"):
            text = extract_pptx_text(file_path)
            type_name = "slides"
        elif ext in (".xlsx", ".xls"):
            text = extract_xlsx_text(file_path)
            type_name = "spreadsheet"
        else:
            print(f"  SKIP (unsupported office format): {file_path}")
            return 0
    except Exception as e:
        print(f"  ERROR extracting {file_path.name}: {e}")
        return 0

    if not text.strip():
        print(f"  SKIP (empty document): {file_path}")
        return 0

    # Split presentations by slide markers, everything else by text chunks
    sections = []
    if type_name == "slides" and "=== SLIDE" in text:
        import re
        slide_splits = re.split(r'===\s*SLIDE\s*\d+\s*===', text)
        sections = [s.strip() for s in slide_splits if s.strip()]
    elif type_name == "spreadsheet" and "=== SHEET" in text:
        import re
        sheet_splits = re.split(r'===\s*SHEET:.*?===', text)
        sections = [s.strip() for s in sheet_splits if s.strip()]
    else:
        sections = chunk_text(
            text,
            chunk_size=config.get("text_chunk_size", DEFAULT_TEXT_CHUNK_SIZE),
            overlap=config.get("text_chunk_overlap", DEFAULT_TEXT_CHUNK_OVERLAP),
        )

    count = 0
    for i, section in enumerate(sections):
        if not section.strip():
            continue
        doc_id = file_id(file_path, i)
        if already_exists(collection, doc_id):
            continue

        embedding = embed_content(client, config, section)
        meta = {
            "source": str(file_path.resolve()),
            "type": type_name,
            "chunk_index": i,
            "total_chunks": len(sections),
            "filename": file_path.name,
            "file_ext": ext,
            "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
        }
        if type_name == "slides":
            meta["slide_number"] = i + 1

        collection.upsert(ids=[doc_id], embeddings=[embedding], documents=[section], metadatas=[meta])
        count += 1
    return count


# Global flag for --force re-ingestion
args_force = False


def ingest_file(client, config, collection, file_path):
    """Route a file to the appropriate ingest handler."""
    file_path = Path(file_path)
    ext = file_path.suffix.lower()

    # Skip common non-content files
    skip_names = {".ds_store", "thumbs.db", ".gitignore", ".gitkeep", "package-lock.json",
                  "yarn.lock", "pnpm-lock.yaml", ".eslintcache"}
    if file_path.name.lower() in skip_names:
        return 0

    # Skip junk directories
    skip_dirs = {".git", "node_modules", "__pycache__", ".venv", "venv", ".env",
                 ".next", ".nuxt", "dist", "build", ".cache", ".turbo",
                 "vendor", ".terraform", ".angular", ".svelte-kit", ".output",
                 "coverage", ".nyc_output", ".pytest_cache", ".mypy_cache"}
    parts = set(file_path.parts)
    if parts & skip_dirs:
        return 0

    # Skip text files > 10MB (likely generated/binary)
    size_mb = file_path.stat().st_size / (1024 * 1024)
    if ext in TEXT_EXTS and size_mb > 10:
        print(f"  SKIP (too large: {size_mb:.0f}MB): {file_path}")
        return 0
    if ext in IMAGE_EXTS and size_mb > 50:
        print(f"  SKIP (too large: {size_mb:.0f}MB): {file_path}")
        return 0
    if ext in DOC_EXTS and size_mb > 100:
        print(f"  SKIP (too large: {size_mb:.0f}MB): {file_path}")
        return 0

    if ext in VIDEO_EXTS:
        return ingest_video(client, config, collection, file_path)
    elif ext in AUDIO_EXTS:
        return ingest_audio(client, config, collection, file_path)
    elif ext in IMAGE_EXTS:
        return ingest_image(client, config, collection, file_path)
    elif ext == ".pdf":
        return ingest_pdf(client, config, collection, file_path)
    elif ext in DOC_EXTS:
        return ingest_office_doc(client, config, collection, file_path)
    elif ext in TEXT_EXTS:
        return ingest_text_file(client, config, collection, file_path)
    else:
        # Try as text for unknown extensions
        try:
            file_path.read_text(errors="strict")[:100]
            return ingest_text_file(client, config, collection, file_path)
        except (UnicodeDecodeError, Exception):
            print(f"  SKIP (binary/unsupported): {file_path}")
            return 0

# ---------------------------------------------------------------------------
# Commands
# ---------------------------------------------------------------------------
def cmd_ingest(args):
    global args_force, _tracker
    args_force = getattr(args, 'force', False)
    _tracker = UsageTracker("ingest")

    config = load_config()
    client = get_genai_client(get_api_key(config))
    collection_name = args.collection or config.get("default_collection", "default")
    collection = get_chroma_collection(collection_name)

    if args_force:
        print(f"Force mode: will re-ingest existing files")

    total = 0
    skipped = 0
    errors = 0

    try:
        for path_str in args.paths:
            p = Path(path_str).resolve()
            if p.is_dir():
                files = sorted(f for f in p.rglob("*") if f.is_file() and not f.name.startswith("."))
                print(f"Ingesting directory: {p} ({len(files)} files)")
                for f in files:
                    print(f"  Processing: {f.relative_to(p)}")
                    try:
                        count = ingest_file(client, config, collection, f)
                        total += count
                        if count > 0:
                            print(f"    Added {count} chunk(s)")
                        elif count == 0:
                            skipped += 1
                    except Exception as e:
                        print(f"    ERROR: {e}")
                        errors += 1
            elif p.is_file():
                print(f"Ingesting: {p.name}")
                try:
                    count = ingest_file(client, config, collection, p)
                    total += count
                    if count > 0:
                        print(f"  Added {count} chunk(s)")
                except Exception as e:
                    print(f"  ERROR: {e}")
                    errors += 1
            else:
                print(f"NOT FOUND: {p}")
    finally:
        _tracker.persist()

    print(f"\nDone! Ingested {total} new chunk(s) into '{collection_name}'")
    if skipped:
        print(f"  Skipped: {skipped} (already existed or empty)")
    if errors:
        print(f"  Errors: {errors}")
    print(_tracker.summary_line())


def deduplicate_results(results, similarity_ratio=0.85):
    """Remove near-duplicate results based on content overlap."""
    if len(results) <= 1:
        return results

    deduped = [results[0]]
    for r in results[1:]:
        is_dup = False
        r_content = r["content"][:500]  # compare first 500 chars
        for existing in deduped:
            e_content = existing["content"][:500]
            # Quick overlap check: count shared words
            r_words = set(r_content.lower().split())
            e_words = set(e_content.lower().split())
            if not r_words or not e_words:
                continue
            overlap = len(r_words & e_words) / max(len(r_words), len(e_words))
            if overlap > similarity_ratio:
                is_dup = True
                break
        if not is_dup:
            deduped.append(r)
    return deduped


def cmd_query(args):
    global _tracker
    _tracker = UsageTracker("query")

    config = load_config()
    client = get_genai_client(get_api_key(config))
    collection_name = args.collection or config.get("default_collection", "default")
    collection = get_chroma_collection(collection_name)

    if collection.count() == 0:
        print("Knowledge base is empty. Ingest some files first.")
        return

    # Fetch extra results so we have room after filtering/dedup
    fetch_k = (args.top_k or 5) * 3
    threshold = args.threshold if args.threshold is not None else config.get("similarity_threshold", DEFAULT_SIMILARITY_THRESHOLD)
    max_tokens = args.max_tokens or config.get("max_tokens", DEFAULT_MAX_TOKENS)
    show_full = args.full
    type_filter = args.type  # e.g., "image", "video", "text", "pdf"

    query_embedding = embed_query(client, config, args.question)

    # Build ChromaDB where filter for type
    where_filter = None
    if type_filter:
        # Map user-friendly names to stored types
        type_map = {
            "image": {"type": "image"},
            "video": {"type": "video_chunk"},
            "text": {"type": "text"},
            "pdf": {"type": "pdf_page"},
            "audio": {"$or": [{"type": "audio"}, {"type": "audio_chunk"}]},
        }
        if type_filter in type_map:
            where_filter = type_map[type_filter]
        else:
            # Direct type match
            where_filter = {"type": type_filter}

    query_kwargs = {
        "query_embeddings": [query_embedding],
        "n_results": min(fetch_k, collection.count()),
        "include": ["documents", "metadatas", "distances"],
    }
    if where_filter:
        query_kwargs["where"] = where_filter

    try:
        results = collection.query(**query_kwargs)
    except Exception as e:
        # where filter might fail if no docs match - fall back to unfiltered
        if where_filter:
            del query_kwargs["where"]
            results = collection.query(**query_kwargs)
        else:
            raise

    # Filter by similarity threshold
    filtered = []
    if results["ids"] and results["ids"][0]:
        for i, doc_id in enumerate(results["ids"][0]):
            distance = results["distances"][0][i] if results["distances"] else 0
            similarity = 1 - distance
            if similarity >= threshold:
                filtered.append({
                    "id": doc_id,
                    "content": results["documents"][0][i] if results["documents"] else "",
                    "similarity": similarity,
                    "metadata": results["metadatas"][0][i] if results["metadatas"] else {},
                })

    # Deduplicate near-identical results (same file in multiple lesson folders)
    filtered = deduplicate_results(filtered)

    # Trim to requested top_k after dedup
    final_k = args.top_k or 5
    filtered = filtered[:final_k]

    # Apply max_tokens budget
    if max_tokens > 0:
        budgeted = []
        token_count = 0
        for r in filtered:
            chunk_tokens = len(r["content"]) // 4
            if token_count + chunk_tokens > max_tokens:
                remaining = max_tokens - token_count
                if remaining > 50:
                    r["content"] = r["content"][:remaining * 4] + "... [truncated]"
                    budgeted.append(r)
                break
            budgeted.append(r)
            token_count += chunk_tokens
        filtered = budgeted

    # Collect unique source files for the agent
    source_files = list(dict.fromkeys(
        r["metadata"].get("source", "") for r in filtered if r["metadata"].get("source")
    ))

    if args.json:
        output = {
            "query": args.question,
            "collection": collection_name,
            "result_count": len(filtered),
            "source_files": source_files,
            "results": [],
        }
        for r in filtered:
            meta = r["metadata"]
            entry = {
                "content": r["content"] if show_full else r["content"][:DEFAULT_PREVIEW_CHARS],
                "content_full_length": len(r["content"]),
                "similarity": round(r["similarity"], 4),
                "source": meta.get("source", ""),
                "type": meta.get("type", ""),
                "filename": meta.get("filename", ""),
            }
            # Add type-specific fields
            if meta.get("chunk_index") is not None:
                entry["chunk_index"] = meta["chunk_index"]
                entry["total_chunks"] = meta.get("total_chunks", 0)
            if meta.get("chunk_start_seconds") is not None:
                entry["time_start"] = meta["chunk_start_seconds"]
                entry["time_end"] = meta.get("chunk_end_seconds", 0)
            if meta.get("chunk_path"):
                entry["chunk_path"] = meta["chunk_path"]
            if meta.get("page_number"):
                entry["page_number"] = meta["page_number"]
            output["results"].append(entry)

        print(json.dumps(output, indent=2))
    else:
        print(f"Query: {args.question}")
        print(f"Collection: {collection_name}")
        print(f"Results: {len(filtered)} (threshold: {threshold})")
        if source_files:
            print(f"Source files ({len(source_files)}):")
            for sf in source_files:
                print(f"  - {sf}")
        print("-" * 60)

        if filtered:
            for i, r in enumerate(filtered):
                meta = r["metadata"]
                print(f"\n[{i+1}] Similarity: {r['similarity']:.3f}")
                print(f"    Source: {meta.get('source', 'unknown')}")
                print(f"    Type: {meta.get('type', 'unknown')}")
                if meta.get("chunk_index") is not None:
                    print(f"    Chunk: {meta['chunk_index'] + 1}/{meta.get('total_chunks', '?')}")
                if meta.get("chunk_start_seconds") is not None:
                    print(f"    Time: {meta['chunk_start_seconds']:.0f}s - {meta.get('chunk_end_seconds', 0):.0f}s")
                if meta.get("chunk_path"):
                    print(f"    Chunk file: {meta['chunk_path']}")
                if meta.get("page_number"):
                    print(f"    Page: {meta['page_number']}")

                content = r["content"]
                if not show_full and len(content) > DEFAULT_PREVIEW_CHARS:
                    content = content[:DEFAULT_PREVIEW_CHARS] + f"... [{len(r['content'])} chars total, use --full to see all]"
                print(f"    Content: {content}")
        else:
            print("No results above similarity threshold.")

    _tracker.persist()


def cmd_usage(args):
    """Show token usage and cost summary."""
    if args.reset:
        if USAGE_FILE.exists():
            USAGE_FILE.unlink()
        print("Usage data reset.")
        return

    if not USAGE_FILE.exists():
        print("No usage data yet. Run an ingest or query first.")
        return

    with open(USAGE_FILE) as f:
        data = json.load(f)

    if args.json:
        print(json.dumps(data, indent=2))
        return

    c = data.get("cumulative", {})
    sessions = data.get("sessions", [])

    emb_cost = (c.get("embedding_tokens", 0) / 1_000_000) * EMBEDDING_PRICE_PER_M
    gen_in_cost = (c.get("generation_input_tokens", 0) / 1_000_000) * FLASH_INPUT_PRICE_PER_M
    gen_out_cost = (c.get("generation_output_tokens", 0) / 1_000_000) * FLASH_OUTPUT_PRICE_PER_M
    total = emb_cost + gen_in_cost + gen_out_cost

    print("mmrag Usage Summary")
    print("=" * 40)
    print(f"Total cost:    ${total:.4f}")
    print(f"Total calls:   {c.get('embedding_calls', 0) + c.get('generation_calls', 0)} "
          f"({c.get('embedding_calls', 0)} embedding, {c.get('generation_calls', 0)} generation)")
    print()
    print("Token breakdown:")
    print(f"  Embedding:         {c.get('embedding_tokens', 0):>10,} tokens (est)  ${emb_cost:.4f}")
    print(f"  Generation input:  {c.get('generation_input_tokens', 0):>10,} tokens        ${gen_in_cost:.4f}")
    print(f"  Generation output: {c.get('generation_output_tokens', 0):>10,} tokens        ${gen_out_cost:.4f}")
    print()
    print(f"Sessions: {len(sessions)}")
    if sessions:
        last = sessions[-1]
        print(f"  Last: {last.get('operation', '?')} @ {last.get('finished_at', '?')} "
              f"(${last.get('cost', {}).get('total', 0):.4f})")

    # Per-day breakdown from sessions
    by_day = {}
    for s in sessions:
        day = s.get("started_at", "")[:10]
        if day:
            by_day.setdefault(day, {"cost": 0, "count": 0})
            by_day[day]["cost"] += s.get("cost", {}).get("total", 0)
            by_day[day]["count"] += 1

    if by_day:
        print()
        print("Daily breakdown:")
        for day in sorted(by_day.keys(), reverse=True)[:7]:
            d = by_day[day]
            print(f"  {day}:  ${d['cost']:.4f} ({d['count']} sessions)")


def cmd_status(args):
    config = load_config()
    collection_name = args.collection or config.get("default_collection", "default")

    try:
        collection = get_chroma_collection(collection_name)
        count = collection.count()
    except Exception:
        count = 0

    print(f"Collection: {collection_name}")
    print(f"Total chunks: {count}")
    print(f"Data dir: {MMRAG_DIR}")
    print(f"ChromaDB: {CHROMADB_DIR}")
    print(f"Config: {CONFIG_FILE}")

    # Show config values
    print(f"\nChunk settings:")
    print(f"  Text: {config.get('text_chunk_size', DEFAULT_TEXT_CHUNK_SIZE)} chars, {config.get('text_chunk_overlap', DEFAULT_TEXT_CHUNK_OVERLAP)} overlap")
    print(f"  Video: {config.get('video_chunk_seconds', DEFAULT_VIDEO_CHUNK_SECONDS)}s, {config.get('video_overlap_seconds', DEFAULT_VIDEO_OVERLAP_SECONDS)}s overlap")
    print(f"  Audio: {config.get('audio_chunk_seconds', DEFAULT_AUDIO_CHUNK_SECONDS)}s, {config.get('audio_overlap_seconds', DEFAULT_AUDIO_OVERLAP_SECONDS)}s overlap")
    print(f"  Embedding dims: {config.get('embedding_dimensions', DEFAULT_EMBEDDING_DIMENSIONS)}")

    if count > 0:
        all_data = collection.get(include=["metadatas"])
        types_map = {}
        sources = set()
        for meta in all_data["metadatas"]:
            t = meta.get("type", "unknown")
            types_map[t] = types_map.get(t, 0) + 1
            sources.add(meta.get("source", "unknown"))

        print(f"\nUnique sources: {len(sources)}")
        print("Type breakdown:")
        for t, c in sorted(types_map.items()):
            print(f"  {t}: {c} chunks")


def cmd_list(args):
    config = load_config()
    collection_name = args.collection or config.get("default_collection", "default")

    try:
        collection = get_chroma_collection(collection_name)
    except Exception:
        print("No data found.")
        return

    all_data = collection.get(include=["metadatas"])
    if not all_data["ids"]:
        print("No documents in collection.")
        return

    # Group by source
    by_source = {}
    for meta in all_data["metadatas"]:
        src = meta.get("source", "unknown")
        if src not in by_source:
            by_source[src] = {"type": meta.get("type", "unknown"), "chunks": 0,
                              "filename": meta.get("filename", "")}
        by_source[src]["chunks"] += 1

    print(f"Collection: {collection_name} ({len(by_source)} files, {len(all_data['ids'])} chunks)")
    print(f"{'Source':<60} {'Type':<15} {'Chunks':<8}")
    print("-" * 85)
    for src, info in sorted(by_source.items()):
        display = src if len(src) <= 58 else "..." + src[-55:]
        print(f"{display:<60} {info['type']:<15} {info['chunks']:<8}")


def cmd_collections(args):
    chroma = get_chroma_client()
    collections = chroma.list_collections()
    if not collections:
        print("No collections found.")
        return
    print(f"{'Collection':<30} {'Documents':<12}")
    print("-" * 44)
    for c in collections:
        col = chroma.get_collection(c.name if hasattr(c, 'name') else c)
        name = c.name if hasattr(c, 'name') else c
        print(f"{name:<30} {col.count():<12}")


def cmd_delete(args):
    config = load_config()
    collection_name = args.collection or config.get("default_collection", "default")
    collection = get_chroma_collection(collection_name)

    source_path = str(Path(args.path).resolve())
    all_data = collection.get(include=["metadatas"])

    ids_to_delete = []
    for i, meta in enumerate(all_data["metadatas"]):
        if meta.get("source") == source_path:
            ids_to_delete.append(all_data["ids"][i])

    if not ids_to_delete:
        print(f"No documents found for: {source_path}")
        return

    collection.delete(ids=ids_to_delete)
    print(f"Deleted {len(ids_to_delete)} chunk(s) from '{collection_name}' for: {source_path}")


def cmd_reset(args):
    if not args.confirm:
        print("ERROR: Pass --confirm to reset the knowledge base.")
        return

    chroma = get_chroma_client()
    collections = chroma.list_collections()
    count = 0
    for c in collections:
        name = c.name if hasattr(c, 'name') else c
        chroma.delete_collection(name)
        count += 1

    import shutil
    if MEDIA_DIR.exists():
        shutil.rmtree(MEDIA_DIR)
        MEDIA_DIR.mkdir(parents=True)

    print(f"Reset complete. Deleted {count} collection(s) and cleared media cache.")

# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser(
        description="Multimodal RAG Knowledge Base CLI",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    sub = parser.add_subparsers(dest="command", help="Command to run")

    # ingest
    p_ingest = sub.add_parser("ingest", help="Ingest files into the knowledge base")
    p_ingest.add_argument("paths", nargs="+", help="File or directory paths to ingest")
    p_ingest.add_argument("--collection", "-c", help="Collection name (default: 'default')")
    p_ingest.add_argument("--force", action="store_true", help="Re-ingest files even if already in the KB")

    # query
    p_query = sub.add_parser("query", help="Query the knowledge base")
    p_query.add_argument("question", help="Question to ask")
    p_query.add_argument("--top-k", "-k", type=int, default=5, help="Max number of results (default: 5)")
    p_query.add_argument("--threshold", "-t", type=float, default=None,
                         help="Min similarity threshold 0.0-1.0 (default: 0.0, return all)")
    p_query.add_argument("--max-tokens", "-m", type=int, default=0,
                         help="Max total tokens in results (0=unlimited)")
    p_query.add_argument("--collection", "-c", help="Collection name")
    p_query.add_argument("--type", help="Filter by content type: image, video, text, pdf, audio")
    p_query.add_argument("--json", "-j", action="store_true", help="Output as JSON (for agent consumption)")
    p_query.add_argument("--full", "-f", action="store_true", help="Show full content (not truncated)")

    # status
    p_status = sub.add_parser("status", help="Show knowledge base status")
    p_status.add_argument("--collection", "-c", help="Collection name")

    # list
    p_list = sub.add_parser("list", help="List ingested documents")
    p_list.add_argument("--collection", "-c", help="Collection name")

    # collections
    sub.add_parser("collections", help="List all collections")

    # delete
    p_delete = sub.add_parser("delete", help="Delete a document by source path")
    p_delete.add_argument("path", help="Source file path to delete")
    p_delete.add_argument("--collection", "-c", help="Collection name")

    # reset
    p_reset = sub.add_parser("reset", help="Reset the entire knowledge base")
    p_reset.add_argument("--confirm", action="store_true", help="Confirm reset")

    # usage
    p_usage = sub.add_parser("usage", help="Show token usage and cost summary")
    p_usage.add_argument("--json", "-j", action="store_true", help="Output as JSON")
    p_usage.add_argument("--reset", action="store_true", help="Reset usage data")

    args = parser.parse_args()

    if not args.command:
        parser.print_help()
        sys.exit(1)

    commands = {
        "ingest": cmd_ingest,
        "query": cmd_query,
        "status": cmd_status,
        "list": cmd_list,
        "collections": cmd_collections,
        "delete": cmd_delete,
        "reset": cmd_reset,
        "usage": cmd_usage,
    }

    commands[args.command](args)


if __name__ == "__main__":
    main()
